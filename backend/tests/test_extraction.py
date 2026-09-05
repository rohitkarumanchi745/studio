"""Microsoft 365 / Graph extraction — auth, encryption-at-rest, ACL→scope
fail-closed retrieval, onboard/delta sync, webhook validation, and graceful
dormancy — proven end to end WITHOUT touching the live network.

Every Graph call is served by an in-process FakeGraph (or a monkeypatched
requests.post for the raw token exchanges), so nothing here reaches Azure. The
background ticker is disabled (STUDIO_GRAPH_SYNC_TICKER=0) so runs are driven
deterministically.

The non-negotiables locked in here:
  - dormancy: with no AZURE_* env the whole feature is inert — imports never
    fail (even with python-docx / python-pptx absent), routes return
    {"configured": false} 200, the ticker starts no thread, onboarding is a
    no-op;
  - the auth abstraction: delegated silent refresh (exactly one exchange,
    re-persisted ciphertext) and app client_credentials caching, both behind one
    interface; for_user() None when dormant / not connected;
  - tokens are Fernet-at-rest, fail closed on a rotated secret, and are NEVER in
    a status body or a raised GraphError;
  - ACL→scope is fail-closed BY CONSTRUCTION: an item outside the user's ACL is
    never ingested, an unresolved ACL is skipped, and a user can retrieve ONLY
    their own private M365 chunks (nobody else's -- not another user and not
    even a Studio admin, whose role wildcard does NOT reach a private 'u:' scope);
  - onboard/delta idempotency: added items ingest once at scope u:{id}, a
    tombstoned item is delete_source'd, a changed item replaces its stable
    source (no duplicate), and the delta cursor advances;
  - the webhook does NO Graph I/O and only nudges next_run_at on a valid,
    constant-time-matched clientState.

Run from the backend directory:
    python -m pytest tests/test_extraction.py -q
"""
import importlib
import io
import sys
import threading
import time
import types
import warnings

import pytest
from fastapi.testclient import TestClient

warnings.filterwarnings("ignore")

_AZURE_KEYS = ("AZURE_TENANT_ID", "AZURE_CLIENT_ID", "AZURE_CLIENT_SECRET")


# ── Env helpers ──────────────────────────────────────────────────────────

def _clear_azure(mp):
    for k in _AZURE_KEYS + ("STUDIO_PUBLIC_URL", "STUDIO_GRAPH_REDIRECT_URI"):
        mp.delenv(k, raising=False)


def _configure_azure(mp):
    mp.setenv("AZURE_TENANT_ID", "tenant-abc")
    mp.setenv("AZURE_CLIENT_ID", "client-abc")
    mp.setenv("AZURE_CLIENT_SECRET", "secret-xyz")


# ── Fixtures ─────────────────────────────────────────────────────────────

@pytest.fixture()
def env(tmp_path, monkeypatch):
    """A fresh SQLite DB with the KAG + Graph tables created, no TestClient.
    configured() reads the env live, so a test toggles Azure per case."""
    monkeypatch.setenv("STUDIO_DB_PATH", str(tmp_path / "m365.db"))
    monkeypatch.setenv("STUDIO_SECRET", "unit-secret-alpha")
    monkeypatch.setenv("STUDIO_GRAPH_SYNC_TICKER", "0")
    for k in ("ANTHROPIC_API_KEY", "OPENAI_API_KEY", "STUDIO_LLM", "STUDIO_LLM_BASE_URL"):
        monkeypatch.delenv(k, raising=False)
    _clear_azure(monkeypatch)

    from app import db as _db
    importlib.reload(_db)                       # repoint DB_PATH at this tmp db
    from app import kag, rbac
    from app.extraction import (acl, client as gclient, graphauth, parsers,
                                store, sync)
    import app.extraction as extraction
    _db.init_db()                              # core tables (users, audit_log, …)
    kag.init_tables()
    store.init_tables()
    parsers.register()
    graphauth._app_token_cache.clear()

    ns = types.SimpleNamespace(
        db=_db, kag=kag, rbac=rbac, acl=acl, gclient=gclient, graphauth=graphauth,
        parsers=parsers, store=store, sync=sync, extraction=extraction,
        mp=monkeypatch, tmp=tmp_path)
    return ns


@pytest.fixture()
def client(tmp_path, monkeypatch):
    """A full TestClient (dormant by default: no AZURE_*). Reloads db then main
    so the app and this test share one fresh db file."""
    monkeypatch.setenv("STUDIO_DB_PATH", str(tmp_path / "m365_app.db"))
    monkeypatch.setenv("STUDIO_SECRET", "route-secret-beta")
    monkeypatch.setenv("STUDIO_GRAPH_SYNC_TICKER", "0")
    monkeypatch.setenv("STUDIO_AUTOPILOT_TICKER", "0")
    for k in ("ANTHROPIC_API_KEY", "OPENAI_API_KEY", "STUDIO_LLM", "STUDIO_LLM_BASE_URL"):
        monkeypatch.delenv(k, raising=False)
    _clear_azure(monkeypatch)

    from app import db as _db
    importlib.reload(_db)
    import app.main as main
    importlib.reload(main)
    with TestClient(main.app) as c:
        yield c


def _login(c, email="admin@studio.local", password="admin123"):
    tok = c.post("/api/auth/login",
                 json={"email": email, "password": password}).json()["access_token"]
    c.headers.update({"Authorization": f"Bearer {tok}"})
    return c


# ── Tiny document builders (in-memory; no assets on disk) ────────────────

def _docx(text):
    from docx import Document
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx(text):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _eml(subject, body):
    from email.message import EmailMessage
    m = EmailMessage()
    m["Subject"] = subject
    m["From"] = "sender@example.com"
    m["To"] = "me@example.com"
    m.set_content(body)
    return bytes(m)


# ── Fake Graph (client + auth) — no network ──────────────────────────────

class FakeGraph:
    """Stands in for GraphClient: same get/get_binary/post/delete surface, plus
    a call log so a test can assert exactly what was (or was NOT) fetched."""

    def __init__(self, drive_feeds=None, perms=None, binaries=None,
                 memberof=("g1",), mail_feed=None):
        self.auth = None                       # bound by FakeAuth
        self.drive_feeds = drive_feeds or {}   # url-substring -> feed dict
        self.perms = perms or {}               # graph_id -> permissions dict | Exception
        self.binaries = binaries or {}         # graph_id -> bytes
        self.memberof = list(memberof)
        self.mail_feed = mail_feed or {"value": [],
                                       "@odata.deltaLink":
                                       "https://graph/v1.0/MAIL_DELTA_1"}
        self.calls = []
        self.posted = []
        self.deleted = []

    def get(self, path, params=None):
        self.calls.append(path)
        if "/memberOf" in path:
            return {"value": [{"id": g} for g in self.memberof]}
        if "/permissions" in path:
            gid = path.split("/drive/items/")[1].split("/permissions")[0]
            v = self.perms.get(gid, {"value": []})
            if isinstance(v, Exception):
                raise v
            return v
        if "messages/delta" in path or "MAIL_DELTA" in path:
            return self.mail_feed
        for key, feed in self.drive_feeds.items():
            if key in path:
                return feed
        return {"value": []}

    def get_binary(self, path):
        self.calls.append(path)
        if "/drive/items/" in path:
            gid = path.split("/drive/items/")[1].split("/content")[0]
            return self.binaries.get(gid, b"")
        if "/messages/" in path:
            gid = path.split("/messages/")[1].split("/$value")[0]
            return self.binaries.get(gid, b"")
        return b""

    def post(self, path, json=None):
        self.posted.append((path, json))
        return {"id": "sub-fake-1"}

    def delete(self, path):
        self.deleted.append(path)


class FakeAuth:
    def __init__(self, user_id, graph_user_id, graph):
        self.user_id = user_id
        self._pid = graph_user_id
        self._graph = graph
        graph.auth = self

    def client(self):
        return self._graph

    def root(self):
        return f"/users/{self._pid}"

    def principal_id(self):
        return self._pid

    def access_token(self):
        return "fake-bearer-never-real"


def _drive_item(gid, name, size=2048, etag="etag-1"):
    return {"id": gid, "name": name, "size": size, "eTag": etag,
            "file": {"mimeType": "application/octet-stream"}}


def _perm_group(gid):
    return {"value": [{"grantedToV2": {"group": {"id": gid}}}]}


def _perm_user(uid):
    return {"value": [{"grantedToV2": {"user": {"id": uid}}}]}


# ── Small raw-requests fake for the token exchanges ──────────────────────

class _FakeResp:
    def __init__(self, ok=True, status_code=200, payload=None):
        self.ok = ok
        self.status_code = status_code
        self._payload = payload or {}

    def json(self):
        return self._payload


# =========================================================================
# Dormancy — the master gate
# =========================================================================

def test_configured_false_without_azure(env):
    assert env.extraction.configured() is False


def test_every_submodule_imports_without_azure_or_parser_wheels(env, monkeypatch):
    # Every submodule imports cleanly with Azure unset (proven at collection too).
    for name in ("app.extraction", "app.extraction.store", "app.extraction.graphauth",
                 "app.extraction.client", "app.extraction.acl",
                 "app.extraction.parsers", "app.extraction.sync",
                 "app.extraction.routes"):
        assert importlib.import_module(name) is not None
    # The ONLY optional wheels are python-docx / python-pptx, imported lazily in
    # parsers.py. Simulate them being absent (importing them raises) and re-exec
    # the module body: it must NOT import them at load, and register() must still
    # succeed. (Only parsers is reloaded — it has no module-level cross-refs, so
    # this leaks no state into the other submodules.)
    monkeypatch.setitem(sys.modules, "docx", None)
    monkeypatch.setitem(sys.modules, "pptx", None)
    reloaded = importlib.reload(importlib.import_module("app.extraction.parsers"))
    reloaded.register()


def test_enqueue_onboard_is_noop_when_dormant(env):
    env.sync.enqueue_onboard("user-dormant")
    assert env.store.get_account("user-dormant") is None


def test_start_ticker_spawns_no_thread_when_dormant(env, monkeypatch):
    # Even with the kill-switch OFF, configured()==False must keep the thread down.
    monkeypatch.setenv("STUDIO_GRAPH_SYNC_TICKER", "1")
    _clear_azure(monkeypatch)
    try:
        env.sync.start_ticker()
        names = {t.name for t in threading.enumerate()}
        assert "m365-sync-ticker" not in names
    finally:
        env.sync.stop_ticker()


def test_status_route_dormant_returns_configured_false(client):
    _login(client)
    for prefix in ("/api",):
        r = client.get(prefix + "/m365/status")
        assert r.status_code == 200
        assert r.json() == {"configured": False}


def test_sync_route_dormant_returns_configured_false_not_500(client):
    _login(client)
    r = client.post("/api/m365/sync")
    assert r.status_code == 200 and r.json() == {"configured": False}


def test_for_user_none_when_dormant(env):
    assert env.graphauth.for_user("whoever") is None


# =========================================================================
# Auth abstraction — one interface, two impls
# =========================================================================

def _account_ct(env, uid):
    c = env.db._conn()
    r = c.execute("SELECT access_ct, refresh_ct FROM graph_accounts WHERE user_id=?",
                  (uid,)).fetchone()
    c.close()
    return (r["access_ct"], r["refresh_ct"]) if r else (None, None)


def test_delegated_refresh_exchanges_once_and_repersists(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-refresh"
    env.store.set_tokens(uid, "OLD-ACCESS", "OLD-REFRESH", time.time() - 60)  # stale
    old_ct = _account_ct(env, uid)

    calls = {"n": 0}

    def fake_post(url, data=None, timeout=None):
        calls["n"] += 1
        assert data["grant_type"] == "refresh_token"
        return _FakeResp(payload={"access_token": "NEW-ACCESS",
                                  "refresh_token": "NEW-REFRESH", "expires_in": 3600})

    import requests
    monkeypatch.setattr(requests, "post", fake_post)

    account = env.store.get_account(uid)
    auth = env.graphauth.DelegatedGraphAuth(uid, account)
    tok = auth.access_token()

    assert tok == "NEW-ACCESS"
    assert calls["n"] == 1                       # exactly one exchange
    new_ct = _account_ct(env, uid)
    assert new_ct != old_ct                      # re-encrypted at rest
    # A second call is fresh now → no further network.
    auth.access_token()
    assert calls["n"] == 1


def test_delegated_valid_token_does_not_refresh(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-fresh"
    env.store.set_tokens(uid, "GOOD-ACCESS", "GOOD-REFRESH", time.time() + 3600)

    import requests
    monkeypatch.setattr(requests, "post",
                        lambda *a, **k: pytest.fail("must not hit the token endpoint"))
    auth = env.graphauth.DelegatedGraphAuth(uid, env.store.get_account(uid))
    assert auth.access_token() == "GOOD-ACCESS"


def test_refresh_failure_raises_grapherror_without_leaking_token(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-badref"
    env.store.set_tokens(uid, "A", "SUPER-SECRET-REFRESH", time.time() - 5)

    import requests
    monkeypatch.setattr(requests, "post",
                        lambda *a, **k: _FakeResp(ok=False, status_code=400))
    auth = env.graphauth.DelegatedGraphAuth(uid, env.store.get_account(uid))
    with pytest.raises(env.gclient.GraphError) as ei:
        auth.access_token()
    assert "SUPER-SECRET-REFRESH" not in str(ei.value)


def test_app_auth_uses_client_credentials_and_caches(env, monkeypatch):
    _configure_azure(monkeypatch)
    env.graphauth._app_token_cache.clear()
    uid = "u-app"
    env.store.save_account(uid, "app", "graph-obj-1")

    calls = {"n": 0}

    def fake_post(url, data=None, timeout=None):
        calls["n"] += 1
        assert data["grant_type"] == "client_credentials"
        assert data["scope"].endswith("/.default")
        return _FakeResp(payload={"access_token": "APP-TOKEN", "expires_in": 3600})

    import requests
    monkeypatch.setattr(requests, "post", fake_post)

    auth = env.graphauth.AppGraphAuth(uid, env.store.get_account(uid))
    assert auth.access_token() == "APP-TOKEN"
    assert auth.access_token() == "APP-TOKEN"    # served from cache
    assert calls["n"] == 1
    assert auth.root() == "/users/graph-obj-1"


def test_for_user_picks_impl_and_none_when_no_row(env, monkeypatch):
    _configure_azure(monkeypatch)
    assert env.graphauth.for_user("nobody") is None    # configured, but no account
    env.store.save_account("u-del", "delegated", "g")
    env.store.save_account("u-app", "app", "g")
    assert isinstance(env.graphauth.for_user("u-del"), env.graphauth.DelegatedGraphAuth)
    assert isinstance(env.graphauth.for_user("u-app"), env.graphauth.AppGraphAuth)


# =========================================================================
# Token storage — Fernet at rest, fail closed
# =========================================================================

def test_token_roundtrip_and_ciphertext_at_rest(env):
    uid = "u-store"
    env.store.set_tokens(uid, "acc-plain", "ref-plain", time.time() + 100)
    got = env.store.get_tokens(uid)
    assert got["access"] == "acc-plain" and got["refresh"] == "ref-plain"
    # Nothing plaintext is persisted.
    ac, rc = _account_ct(env, uid)
    assert "acc-plain" not in ac and "ref-plain" not in rc


def test_corrupted_ciphertext_fails_closed_to_none(env):
    uid = "u-corrupt"
    env.store.set_tokens(uid, "acc", "ref", time.time() + 100)
    c = env.db._conn()
    c.execute("UPDATE graph_accounts SET access_ct=? WHERE user_id=?", ("garbage!!", uid))
    c.commit()
    c.close()
    assert env.store.get_tokens(uid) is None


def test_rotated_secret_revokes_account_and_never_leaks(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-rotate"
    env.store.set_tokens(uid, "PLAINTEXT-ACCESS", "PLAINTEXT-REFRESH", time.time() + 3600)
    # Rotate the master secret: stored tokens become undecryptable.
    monkeypatch.setenv("STUDIO_SECRET", "a-completely-different-secret")
    assert env.store.get_tokens(uid) is None
    auth = env.graphauth.DelegatedGraphAuth(uid, env.store.get_account(uid))
    with pytest.raises(env.gclient.GraphError) as ei:
        auth.access_token()
    assert "PLAINTEXT" not in str(ei.value)
    assert env.store.get_account(uid)["status"] == "revoked"


def test_public_status_exposes_no_ciphertext(env):
    uid = "u-pub"
    env.store.save_account(uid, "delegated", "g")
    env.store.set_tokens(uid, "acc", "ref", time.time() + 100)
    ps = env.store.public_status(uid)
    blob = repr(ps)
    ac, rc = _account_ct(env, uid)
    assert ac not in blob and rc not in blob
    assert not any(k in ps for k in ("access", "refresh", "access_ct", "refresh_ct"))


# =========================================================================
# ACL → scope — fail closed BY CONSTRUCTION
# =========================================================================

class _TinyClient:
    """Just enough of GraphClient for acl.py: .auth.root() + .get()."""

    def __init__(self, perms, memberof):
        self.auth = types.SimpleNamespace(root=lambda: "/me")
        self._perms = perms
        self._memberof = memberof

    def get(self, path, params=None):
        if "/memberOf" in path:
            return {"value": [{"id": g} for g in self._memberof]}
        if "/permissions" in path:
            gid = path.split("/drive/items/")[1].split("/permissions")[0]
            v = self._perms.get(gid)
            if isinstance(v, Exception):
                raise v
            return v or {"value": []}
        return {}


def test_principal_set_is_object_id_union_groups(env):
    cli = _TinyClient({}, memberof=["grp-a", "grp-b"])
    ps = env.acl.principal_set(cli, "obj-1")
    assert ps == {"obj-1", "grp-a", "grp-b"}


def test_principal_set_empty_on_failure(env):
    class Boom(_TinyClient):
        def get(self, path, params=None):
            raise RuntimeError("throttled")
    assert env.acl.principal_set(Boom({}, []), "obj-1") == set()


def test_resolve_grants_on_group_intersection(env):
    cli = _TinyClient({"itm": _perm_group("grp-a")}, memberof=["grp-a"])
    principals = env.acl.principal_set(cli, "obj-1")
    scope = env.acl.resolve_access_scope({"id": "u9"}, "drive", "itm", cli, principals)
    assert scope == "u:u9"


def test_resolve_denies_when_no_intersection(env):
    cli = _TinyClient({"itm": _perm_user("someone-else")}, memberof=["grp-a"])
    principals = env.acl.principal_set(cli, "obj-1")
    assert env.acl.resolve_access_scope({"id": "u9"}, "drive", "itm", cli, principals) is None


def test_resolve_denies_when_permissions_call_raises(env):
    cli = _TinyClient({"itm": RuntimeError("500")}, memberof=["grp-a"])
    principals = env.acl.principal_set(cli, "obj-1")
    assert env.acl.resolve_access_scope({"id": "u9"}, "drive", "itm", cli, principals) is None


def test_resolve_denies_on_empty_principal_set(env):
    cli = _TinyClient({"itm": _perm_group("grp-a")}, memberof=["grp-a"])
    assert env.acl.resolve_access_scope({"id": "u9"}, "drive", "itm", cli, set()) is None


def test_mail_owner_always_granted(env):
    cli = _TinyClient({}, memberof=[])
    assert env.acl.resolve_access_scope({"id": "u9"}, "mail", "msg", cli, set()) == "u:u9"


def test_anonymous_link_does_not_widen(env):
    # A link/anonymous grant carries no user/group id → no intersection → deny.
    perms = {"itm": {"value": [{"link": {"scope": "anonymous"}}]}}
    cli = _TinyClient(perms, memberof=["grp-a"])
    principals = env.acl.principal_set(cli, "obj-1")
    assert env.acl.resolve_access_scope({"id": "u9"}, "drive", "itm", cli, principals) is None


# =========================================================================
# rbac / kag threading — additive, existing behaviour byte-identical
# =========================================================================

def test_scope_sql_has_no_percent_literal_postgres_safe(env):
    # _pg_sql rewrites ? -> %s for psycopg, so a literal % in the SQL breaks
    # parameter binding on Postgres (but not sqlite). Guard against reintroducing
    # a LIKE 'u:%' style predicate that only fails in prod.
    for role, uid in [("admin", "u9"), ("admin", None), ("analyst", "u9"),
                      ("viewer", None), ("admin", "adm")]:
        pred, params = env.kag._scope_sql(role, uid)
        assert "%" not in pred, f"scope SQL for {role}/{uid} contains a literal %: {pred}"


def test_kag_scopes_for_no_user_id_is_unchanged(env):
    assert env.rbac.kag_scopes_for("viewer") == {"viewer"}
    assert env.rbac.kag_scopes_for("analyst") == {"analyst"}
    assert env.rbac.kag_scopes_for("admin") == "*"


def test_kag_scopes_for_adds_only_own_private_scope(env):
    assert env.rbac.kag_scopes_for("viewer", "u9") == {"viewer", "u:u9"}
    assert env.rbac.kag_scopes_for("admin", "u9") == "*"        # admin unchanged


# =========================================================================
# Parsers — docx/pptx registration, sanitizer, oversize
# =========================================================================

def test_docx_and_pptx_parse_through_kag_after_register(env):
    env.parsers.register()
    st, chunks = env.kag.parse("Deck.docx", _docx("Acme quarterly revenue projections."))
    assert st == "docx" and chunks and chunks[0]["text"].strip()
    st2, chunks2 = env.kag.parse("Deck.pptx", _pptx("Roadmap for the next quarter"))
    assert st2 == "pptx" and chunks2


def test_sanitize_name_coerces_to_name_re(env):
    out = env.parsers.sanitize_name("Q3/Report:final – v2 (draft).docx")
    assert env.kag._NAME_RE.match(out) and len(out) <= 64
    long_unicode = "é" * 200
    out2 = env.parsers.sanitize_name(long_unicode, "fallback")
    assert env.kag._NAME_RE.match(out2) and len(out2) <= 64


def test_stable_source_name_is_deterministic_and_valid(env):
    # The same item (same id + name) re-synced maps to the SAME source name, so a
    # re-ingest is a wholesale replace rather than a duplicate.
    a = env.parsers.stable_source_name("graph-id-123", "Weird/Name:!.docx")
    b = env.parsers.stable_source_name("graph-id-123", "Weird/Name:!.docx")
    assert a == b
    # The id-derived hash tail distinguishes different items even with the same
    # human name, so two files never collide on one source name.
    c = env.parsers.stable_source_name("graph-id-999", "Weird/Name:!.docx")
    assert c != a
    assert env.kag._NAME_RE.match(a) and a.endswith(".docx") and len(a) <= 64


def test_oversize_item_skipped_without_raising(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-big"
    env.store.save_account(uid, "app", "obj-big")
    env.sync._ensure_collection(uid)
    graph = FakeGraph(
        drive_feeds={"/drive/root/delta": {
            "value": [_drive_item("big", "Huge.docx",
                                  size=env.sync.MAX_ITEM_BYTES + 1)],
            "@odata.deltaLink": "https://graph/v1.0/DRIVE_DELTA_1"}},
        perms={"big": _perm_group("g1")}, binaries={"big": b"x"})
    monkeypatch.setattr(env.graphauth, "for_user",
                        lambda u: FakeAuth(uid, "obj-big", graph))
    stats = env.sync.run_sync(uid)
    assert stats["ok"] and stats["ingested"] == 0 and stats["skipped"] >= 1


# =========================================================================
# Sync — onboard, ACL-scoped ingest, per-user isolation, idempotency, tombstone
# =========================================================================

def _wire_run(env, monkeypatch, uid, graph, graph_user_id="obj-1"):
    env.store.save_account(uid, "app", graph_user_id)
    env.sync._ensure_collection(uid)
    monkeypatch.setattr(env.graphauth, "for_user",
                        lambda u: FakeAuth(uid, graph_user_id, graph))


def test_acl_denied_item_never_ingested_and_isolation(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "owner-1"
    ok_bytes = _docx("Acme quarterly revenue projections and margins.")
    graph = FakeGraph(
        drive_feeds={"/drive/root/delta": {
            "value": [_drive_item("ok", "Allowed.docx"),
                      _drive_item("no", "Denied.docx", etag="etag-2")],
            "@odata.deltaLink": "https://graph/v1.0/DRIVE_DELTA_1"}},
        perms={"ok": _perm_group("g1"), "no": _perm_user("stranger")},
        binaries={"ok": ok_bytes, "no": _docx("secret you cannot see")})
    _wire_run(env, monkeypatch, uid, graph)

    stats = env.sync.run_sync(uid)
    assert stats["ok"] and stats["ingested"] == 1 and stats["skipped"] >= 1

    # The denied doc produced NO chunk anywhere.
    c = env.db._conn()
    n_denied = c.execute(
        "SELECT COUNT(*) AS n FROM kag_chunks WHERE chunk_text LIKE ?",
        ("%secret you cannot see%",)).fetchone()["n"]
    c.close()
    assert n_denied == 0

    # Per-user isolation at retrieval: ONLY the owner reaches it. Not another
    # user, and NOT even a Studio admin (whose role wildcard does not reach a
    # private 'u:' scope) -- the Microsoft source ACL is the authority, not
    # Studio's role hierarchy. An admin who IS the owner still reaches it.
    q = "quarterly revenue"
    assert env.kag.search(q, role="viewer", user_id=uid)                    # owner
    assert env.kag.search(q, role="viewer", user_id="other-user") == []     # another user
    assert env.kag.search(q, role="admin", user_id="admin-x") == []         # admin != owner: denied
    assert env.kag.search(q, role="admin", user_id=uid)                     # admin who IS the owner

    # And the same isolation on collection LISTING (no metadata leak of which
    # users have M365 connected): a non-owner admin's reachable set excludes it.
    names = lambda r, u=None: {c["name"] for c in env.kag.reachable_collections(r, u)}
    assert f"m365-{uid}" in names("viewer", uid)
    assert f"m365-{uid}" not in names("admin", "admin-x")
    assert f"m365-{uid}" in names("admin", uid)


def test_onboard_delta_idempotency_and_tombstone(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "owner-2"
    a_v1 = _docx("Alpha report first version revenue.")
    a_v2 = _docx("Alpha report SECOND version revenue growth.")
    b_bytes = _docx("Beta memo about hiring.")

    feed1 = {"value": [_drive_item("A", "Alpha.docx", etag="A1"),
                       _drive_item("B", "Beta.docx", etag="B1")],
             "@odata.deltaLink": "https://graph/v1.0/DRIVE_DELTA_1"}
    feed2 = {"value": [_drive_item("A", "Alpha.docx", etag="A2"),            # changed
                       {"id": "B", "@removed": {"reason": "deleted"}}],       # tombstone
             "@odata.deltaLink": "https://graph/v1.0/DRIVE_DELTA_2"}
    graph = FakeGraph(
        drive_feeds={"/drive/root/delta": feed1, "DRIVE_DELTA_1": feed2},
        perms={"A": _perm_group("g1"), "B": _perm_group("g1")},
        binaries={"A": a_v1, "B": b_bytes})
    _wire_run(env, monkeypatch, uid, graph)

    # Spy on delete_source while keeping the real implementation.
    deletes = []
    real_delete = env.kag.delete_source

    def spy_delete(collection, source_name):
        deletes.append(source_name)
        return real_delete(collection, source_name)
    monkeypatch.setattr(env.kag, "delete_source", spy_delete)

    s1 = env.sync.run_sync(uid)
    assert s1["ingested"] == 2
    coll = f"m365-{uid}"
    a_src = env.parsers.stable_source_name("A", "Alpha.docx")
    b_src = env.parsers.stable_source_name("B", "Beta.docx")

    def _sources():
        c = env.db._conn()
        rows = c.execute(
            "SELECT DISTINCT source_name FROM kag_chunks WHERE collection=?",
            (coll,)).fetchall()
        c.close()
        return {r["source_name"] for r in rows}

    assert _sources() == {a_src, b_src}
    # Delta cursor was stored → the next run walks from it, not the initial feed.
    assert env.store.get_delta(uid, "drive") == feed1["@odata.deltaLink"]

    # Second run: A changed (same stable source, wholesale replace), B removed.
    graph.binaries["A"] = a_v2
    s2 = env.sync.run_sync(uid)
    assert s2["removed"] == 1
    assert b_src in deletes                       # the tombstoned item was delete_source'd
    assert _sources() == {a_src}                  # no duplicate; B gone

    # A's chunks now reflect the new version only.
    c = env.db._conn()
    txt = " ".join(r["chunk_text"] for r in c.execute(
        "SELECT chunk_text FROM kag_chunks WHERE collection=? AND source_name=?",
        (coll, a_src)).fetchall())
    c.close()
    assert "SECOND version" in txt and "first version" not in txt
    assert env.store.get_delta(uid, "drive") == feed2["@odata.deltaLink"]


def test_enqueue_onboard_when_configured_creates_collection_and_schedules(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-onboard"
    env.sync.enqueue_onboard(uid)
    acct = env.store.get_account(uid)
    assert acct is not None and acct["status"] == "onboarding"
    assert acct["next_run_at"] is not None
    assert env.kag._collection(f"m365-{uid}")["access_scope"] == "u:" + uid


# =========================================================================
# Webhook — no Graph I/O, constant-time clientState
# =========================================================================

def test_process_notification_valid_and_invalid_client_state(env, monkeypatch):
    _configure_azure(monkeypatch)
    uid = "u-hook"
    env.store.save_account(uid, "delegated", "g")
    env.store.update_fields(uid, {"next_run_at": 100.0})
    env.store.save_subscription("sub-1", uid, "/me/drive/root",
                                "the-shared-secret", time.time() + 3600)

    # Assert the webhook path never builds a Graph client.
    monkeypatch.setattr(env.graphauth, "for_user",
                        lambda u: pytest.fail("webhook must not fetch Graph"))

    assert env.sync.process_notification("sub-1", "wrong-secret") is False
    assert env.store.get_account(uid)["next_run_at"] == 100.0    # unchanged

    assert env.sync.process_notification("sub-1", "the-shared-secret") is True
    assert env.store.get_account(uid)["next_run_at"] > 100.0     # nudged


def test_webhook_route_echoes_validation_token(client):
    for prefix in ("/api",):
        r = client.post(prefix + "/m365/webhook?validationToken=hs-123")
        assert r.status_code == 200 and r.text == "hs-123"
        assert r.headers["content-type"].startswith("text/plain")
        rg = client.get(prefix + "/m365/webhook?validationToken=hs-xyz")
        assert rg.status_code == 200 and rg.text == "hs-xyz"


# =========================================================================
# Scheduler claim — race-safe (a claimed row is not re-claimed)
# =========================================================================

def test_claim_due_is_race_safe(env):
    for uid in ("c1", "c2"):
        env.store.save_account(uid, "app", "g")
        env.store.update_fields(uid, {"status": "connected", "next_run_at": 1.0})
    # Distinct ticks (as the real ticker uses time.time()) claim disjoint sets:
    # the first sweep wins both due rows; the second finds them freshly claimed.
    first = {a["user_id"] for a in env.store.claim_due(time.time())}
    second = {a["user_id"] for a in env.store.claim_due(time.time() + 1)}
    assert first == {"c1", "c2"}
    assert first.isdisjoint(second)
    # A revoked account is never claimed.
    env.store.update_fields("c1", {"status": "revoked", "next_run_at": 1.0,
                                   "claimed_at": None})
    assert all(a["user_id"] != "c1"
               for a in env.store.claim_due(time.time() + 1000))


# =========================================================================
# Routes — configured connect (delegated + app), no token in body
# =========================================================================

def test_connect_delegated_returns_authorize_url(client, monkeypatch):
    _configure_azure(monkeypatch)
    monkeypatch.setenv("STUDIO_GRAPH_AUTH_MODE", "delegated")
    _login(client)
    r = client.post("/api/m365/connect")
    assert r.status_code == 200
    url = r.json()["authorize_url"]
    assert "login.microsoftonline.com/tenant-abc" in url
    assert "offline_access" in url and "state=" in url
    assert "access_token" not in url


def test_connect_app_mode_provisions_row(client, monkeypatch):
    _configure_azure(monkeypatch)
    monkeypatch.setenv("STUDIO_GRAPH_AUTH_MODE", "app")
    _login(client)
    r = client.post("/api/m365/connect")
    assert r.status_code == 200 and r.json().get("connected") is True

    from app.extraction import store
    admin = client.get("/api/auth/me").json()
    acct = store.get_account(admin["id"])
    assert acct is not None and acct["auth_mode"] == "app"


def test_status_never_returns_a_token(client, monkeypatch):
    _configure_azure(monkeypatch)
    _login(client)
    admin = client.get("/api/auth/me").json()
    from app.extraction import store
    store.save_account(admin["id"], "delegated", "g")
    store.set_tokens(admin["id"], "TOP-SECRET-ACCESS", "TOP-SECRET-REFRESH",
                     time.time() + 3600)
    r = client.get("/api/m365/status")
    assert r.status_code == 200
    body = r.text
    assert "TOP-SECRET" not in body
    assert "_ct" not in body and "refresh" not in body
    assert r.json()["configured"] is True and r.json()["connected"] is True
