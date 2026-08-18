"""The Word/.docx and PowerPoint/.pptx parsers kag.py was missing (it shipped
only xlsx/csv/pdf/eml), plus the Graph name sanitizer.

register() splices 'docx' and 'pptx' into kag's parser registry AT RUNTIME —
kag._PARSERS, kag._SOURCE_TYPES and kag._sniff_type — so kag.ingest_bytes accepts
those bytes unchanged, with no edit to kag's parse/store/retrieve core. It is
idempotent (guarded by a marker attribute) and re-entrant.

python-docx / python-pptx are imported LAZILY inside the parser functions, so a
missing optional wheel can never break `import app.extraction.parsers` at
startup (dormancy) — it only surfaces as a clean parse error if such a file is
actually ingested.

Every Graph-derived name (filename, mail subject) is UNTRUSTED metadata: it is
coerced through sanitize_name() / stable_source_name() to kag._NAME_RE
(^[A-Za-z0-9 _.-]{1,64}$) before it is ever used as a KAG source name, so a
crafted filename can neither carry SQL/identifier metacharacters nor collide
across items. stable_source_name() keys on the Graph item id (hashed) so a
re-synced/changed item maps to the SAME source name — a wholesale replace, never
a duplicate — while preserving the real extension so kag's sniffer still routes
it to the right parser.
"""
import hashlib
import io
import os
import re

_DISALLOWED = re.compile(r"[^A-Za-z0-9 _.\-]")


def register():
    """Idempotently teach kag about docx/pptx. Safe to call at kag import/init."""
    from .. import kag
    if getattr(kag, "_extraction_parsers_registered", False):
        return
    kag._PARSERS["docx"] = parse_docx
    kag._PARSERS["pptx"] = parse_pptx
    if "docx" not in kag._SOURCE_TYPES:
        kag._SOURCE_TYPES = tuple(kag._SOURCE_TYPES) + ("docx", "pptx")

    # kag._sniff_type maps a PK/zip container to 'xlsx'; a .docx/.pptx is also a
    # PK zip, so route those by extension FIRST, then defer to the original.
    _orig_sniff = kag._sniff_type

    def _sniff(name, data):
        ext = os.path.splitext((name or "").lower())[1].lstrip(".")
        if ext == "docx":
            return "docx"
        if ext == "pptx":
            return "pptx"
        return _orig_sniff(name, data)

    kag._sniff_type = _sniff
    kag._extraction_parsers_registered = True


def _windows(text):
    from .. import kag
    return kag._windows(text)


def parse_docx(data, source_name):
    """Word document → chunk dicts. Body paragraphs and table cells, split into
    kag's char windows. python-docx imported lazily."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    author = None
    try:
        author = doc.core_properties.author
    except Exception:
        pass
    parts = [p.text for p in doc.paragraphs if (p.text or "").strip()]
    for table in getattr(doc, "tables", []):
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = " | ".join(x for x in cells if x)
            if line:
                parts.append(line)
    body = "\n".join(parts)
    chunks = []
    for piece in _windows(body):
        chunks.append({"text": piece, "meta": {"author": author}})
    return chunks


def parse_pptx(data, source_name):
    """PowerPoint → one chunk per slide (long slides windowed). python-pptx
    imported lazily."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    slide_count = len(slides)
    chunks = []
    for idx, slide in enumerate(slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        slide_text = "\n".join(texts).strip()
        if not slide_text:
            continue
        for piece in (_windows(slide_text) or [slide_text]):
            chunks.append({"text": piece, "meta": {"slide": idx, "slide_count": slide_count}})
    return chunks


def sanitize_name(raw, fallback="document"):
    """Coerce an arbitrary Graph filename/subject to a kag._NAME_RE-valid name,
    <=64 chars. Disallowed characters (incl. any non-ASCII) become '_'."""
    cleaned = _DISALLOWED.sub("_", (raw or "").strip())[:64].strip()
    if cleaned and _valid(cleaned):
        return cleaned
    fb = _DISALLOWED.sub("_", (fallback or "").strip())[:64].strip() or "document"
    return fb if _valid(fb) else "document"


def _valid(name):
    from .. import kag
    return bool(kag._NAME_RE.match(name or ""))


def stable_source_name(graph_id, name):
    """A NAME_RE-valid source name that is STABLE across re-sync (keyed on the
    Graph item id) yet keeps a readable stem and the real extension so kag's
    sniffer routes it correctly. Form: '<stem>_<12-hex-of-graph-id>.<ext>'."""
    name = name or "document"
    stem, dot, ext = name.rpartition(".")
    if not dot:
        stem, ext = name, ""
    gid = hashlib.sha256((graph_id or "").encode("utf-8")).hexdigest()[:12]
    ext_clean = re.sub(r"[^A-Za-z0-9]", "", ext.lower())[:8]
    tail = f"_{gid}" + (f".{ext_clean}" if ext_clean else "")
    room = max(1, 64 - len(tail))
    base = sanitize_name(stem, "doc")[:room].strip() or "doc"
    out = f"{base}{tail}"
    if _valid(out):
        return out
    fallback = (gid + (f".{ext_clean}" if ext_clean else ""))[:64]
    return fallback if _valid(fallback) else gid
